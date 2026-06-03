"""Markdown → PowerPoint (.pptx) converter.

One chapter per slide. Tables are converted to simplified bullet lists.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from finance_agent.export.parser import Section, parse_markdown, split_by_chapters


def markdown_to_pptx(markdown_text: str, output_path: str, stock_name: str = "") -> str:
    """Convert markdown report to PowerPoint presentation.

    Parameters
    ----------
    markdown_text : str
        Full markdown report.
    output_path : str
        Destination file path.
    stock_name : str
        Stock name for the cover slide.

    Returns
    -------
    str
        The output_path.
    """
    prs = Presentation()

    # Cover slide
    _add_cover_slide(prs, stock_name or "分析报告")

    # Parse and split into chapters
    sections = parse_markdown(markdown_text)
    chapters = split_by_chapters(sections)

    for chapter_title, chapter_sections in chapters:
        if not chapter_title:
            continue
        _add_chapter_slide(prs, chapter_title, chapter_sections)

    # Save
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _add_cover_slide(prs: Presentation, title_text: str) -> None:  # pyrefly: ignore[not-a-type]
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "AI 投资分析报告"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER


def _add_chapter_slide(prs: Presentation, title: str, sections: list[Section]) -> None:  # pyrefly: ignore[not-a-type]
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Title
    slide.shapes.title.text = title

    # Content
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for sec in sections:
        if sec.type == "heading":
            p = body.add_paragraph()
            p.text = sec.text
            p.level = min(sec.level, 2)
            p.font.bold = True
        elif sec.type == "paragraph":
            p = body.add_paragraph()
            p.text = sec.text
            p.level = 0
            p.font.size = Pt(14)
        elif sec.type == "table":
            # Convert table to simplified bullet list
            if sec.rows:
                header = sec.rows[0]
                for row in sec.rows[1:]:
                    p = body.add_paragraph()
                    cells = [f"{h}: {c}" for h, c in zip(header, row, strict=False)]
                    p.text = " | ".join(cells)
                    p.level = 1
                    p.font.size = Pt(12)
        elif sec.type == "separator":
            p = body.add_paragraph()
            p.text = ""
